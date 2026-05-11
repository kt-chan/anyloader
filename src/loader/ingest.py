import os
import lancedb
import pandas as pd
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from pathlib import Path
import httpx
from dotenv import load_dotenv
import uuid

load_dotenv()

LANCEDB_URI = os.getenv("LANCEDB_URI", "./.lancedb")
TABLE_NAME = "course_prerequisites"
LLM_HOST_URL = os.getenv("LLM_HOST_URL")
LLM_API_KEY = os.getenv("API_KEY") # Fixed from EMBEDDING_API_KEY
MODEL_NAME = os.getenv("MODEL_NAME")
INPUT_DIR = os.getenv("INPUT_DIR", "data")

class Ingestor:
    def __init__(self):
        self.db = lancedb.connect(LANCEDB_URI)
        self.table = None

    async def get_embeddings(self, texts: list[str]):
        # Zhipu AI has a batch limit of 64 items per request
        batch_size = 64
        url = f"{LLM_HOST_URL}/embeddings"
        headers = {
            "Authorization": f"Bearer {LLM_API_KEY}",
            "Content-Type": "application/json"
        }
        
        all_embeddings = []
        
        async with httpx.AsyncClient() as client:
            for i in range(0, len(texts), batch_size):
                batch = [t for t in texts[i:i + batch_size] if t.strip()]
                if not batch:
                    continue
                    
                payload = {
                    "model": "embedding-3",
                    "input": batch
                }
                
                # Retry logic for rate limiting or transient errors
                max_retries = 3
                for attempt in range(max_retries):
                    try:
                        response = await client.post(url, headers=headers, json=payload, timeout=60.0)
                        if response.status_code == 200:
                            data = response.json()
                            all_embeddings.extend([item["embedding"] for item in data["data"]])
                            break
                        elif response.status_code == 429 or (response.status_code == 400 and "1210" in response.text):
                            # Rate limit or suspicious param error (might be congestion)
                            wait_time = (attempt + 1) * 2
                            print(f"Batch {i//batch_size} failed (Status {response.status_code}). Retrying in {wait_time}s... (Attempt {attempt + 1}/{max_retries})")
                            await asyncio.sleep(wait_time)
                        else:
                            print(f"Embedding API failed for batch {i//batch_size} with status {response.status_code}: {response.text}")
                            if attempt == max_retries - 1:
                                all_embeddings.extend([[0.1] * 2048 for _ in batch])
                            await asyncio.sleep(1)
                    except Exception as e:
                        print(f"Error calling embedding API for batch {i//batch_size}: {e}")
                        if attempt == max_retries - 1:
                            all_embeddings.extend([[0.1] * 2048 for _ in batch])
                        await asyncio.sleep(1)
                
                # Small mandatory delay between batches to respect rate limits
                await asyncio.sleep(0.5)

        return all_embeddings

    def extract_text(self, file_path: Path):
        ext = file_path.suffix.lower()
        text = ""
        try:
            if ext == ".pdf":
                with fitz.open(file_path) as doc:
                    for page in doc:
                        text += page.get_text()
            elif ext == ".docx":
                doc = Document(file_path)
                text = "\n".join([para.text for para in doc.paragraphs])
            elif ext == ".pptx":
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            elif ext in [".xlsx", ".xls"]:
                df = pd.read_excel(file_path)
                text = df.to_string()
            elif ext in [".md", ".txt"]:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()
        except Exception as e:
            print(f"Error extracting text from {file_path}: {e}")
        return text

    def chunk_text(self, text: str, chunk_size=1000, overlap=100):
        chunks = []
        for i in range(0, len(text), chunk_size - overlap):
            chunks.append(text[i:i + chunk_size])
        return chunks

    async def process_data_dir(self, data_dir: str):
        import re
        data = []
        path = Path(data_dir)
        year_pattern = re.compile(r"\d{4}-\d{4}")
        
        for file_path in path.rglob("*"):
            if file_path.is_file():
                relative_path = file_path.relative_to(path)
                parts = relative_path.parts
                
                metadata = {
                    "degree_level": "unknown",
                    "category": "unknown",
                    "department": "unknown",
                    "academic_year": "unknown",
                    "faculty": "generic",
                    "doc_title": file_path.name.lower().replace(" ", "_"),
                    "file_path": str(file_path)
                }
                
                # Inference logic
                if len(parts) > 0:
                    metadata["degree_level"] = parts[0].lower().replace(" ", "_")
                
                # Look for academic year and faculty
                year_idx = -1
                for i, part in enumerate(parts):
                    if year_pattern.match(part):
                        metadata["academic_year"] = part.lower().replace(" ", "_")
                        year_idx = i
                        break
                
                if year_idx != -1:
                    # Category is usually between degree_level and academic_year
                    if year_idx > 1:
                        metadata["category"] = parts[1].lower().replace(" ", "_")
                    
                    # Faculty is usually after academic_year
                    if len(parts) > year_idx + 1:
                        metadata["faculty"] = parts[year_idx + 1].lower().replace(" ", "_")
                    
                    # Try to infer department if there's more depth
                    if len(parts) > year_idx + 2:
                        dept = parts[year_idx + 2]
                        metadata["department"] = dept.lower().replace(" ", "_") if not dept.endswith(file_path.suffix) else "unknown"
                elif len(parts) > 2:
                    # Fallback for paths without year
                    metadata["category"] = parts[1].lower().replace(" ", "_")
                
                # Extract year for range filtering
                try:
                    metadata["start_year"] = int(metadata["academic_year"].split("-")[0])
                except:
                    metadata["start_year"] = 0

                content = self.extract_text(file_path)
                chunks = self.chunk_text(content)
                
                for chunk in chunks:
                    chunk = chunk.strip()
                    if chunk:
                        # Heuristic to extract section: use the first line if it's short
                        first_line = chunk.split("\n")[0].strip()
                        section = "general"
                        if 3 < len(first_line) < 60:
                            section = first_line.lower().replace(" ", "_").strip("_")
                        
                        data.append({
                            "text": chunk,
                            "section": section,
                            **metadata
                        })

        if not data:
            print("No data found to ingest.")
            return

        # Get embeddings
        texts = [d["text"] for d in data]
        embeddings = await self.get_embeddings(texts)
        
        for i, d in enumerate(data):
            d["vector"] = embeddings[i]

        df = pd.DataFrame(data)
        # Always overwrite for now to ensure schema is updated with 'department'
        self.table = self.db.create_table(TABLE_NAME, data=df, mode="overwrite")
        
        # Create FTS index for doc_title fuzzy match
        self.table.create_fts_index("text", replace=True) 

if __name__ == "__main__":
    import asyncio
    ingestor = Ingestor()
    asyncio.run(ingestor.process_data_dir(INPUT_DIR))
