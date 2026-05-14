from fastapi import FastAPI, HTTPException, File, UploadFile, Form, BackgroundTasks
from pydantic import BaseModel, Field
from typing import Optional, List, Dict, Any
import lancedb
import os
import httpx
import asyncio
import re
import io
import time
import uuid
import shutil
from pathlib import Path
from PyPDF2 import PdfReader
import pandas as pd
from dotenv import load_dotenv

load_dotenv()

# --- Configuration and Database Initialization ---
LANCEDB_URI = os.getenv("LANCEDB_URI", "./storage/lancedb")
TABLE_NAME = "course_chunks"
RAG_HOST_PATH = os.getenv("RAG_HOST_PATH", "127.0.0.1")
RAG_PORT = int(os.getenv("RAG_PORT", "8000"))
API_KEY = os.getenv("API_KEY")
LLM_HOST_URL = os.getenv("LLM_HOST_URL") # Base URL like https://open.bigmodel.cn/api/paas/v4
EMBEDDING_MODEL_NAME = os.getenv("EMBEDDING_MODEL_NAME", "embedding-2")
LLM_MODEL_NAME = os.getenv("LLM_MODEL_NAME", "glm-4-flash")

TEMP_DIR = Path("storage/temp")
TEMP_DIR.mkdir(parents=True, exist_ok=True)

# Initialize LanceDB
db = lancedb.connect(LANCEDB_URI)

# --- Global Storage (In-memory) ---
uploaded_files_store = {}  # {file_id: {"name": ..., "content": ..., "metadata": ...}}
tasks_store = {}           # {task_id: {"status": ..., "error": ...}}
file_id_counter = 0

# --- FastAPI App Initialization ---
app = FastAPI(title="anyLoader RAG API")

# --- Pydantic Models ---
class FileMetadata(BaseModel):
    id: int
    name: str
    size: int
    extension: str
    mime_type: str
    created_at: int

class ProcessRequest(BaseModel):
    file_id: int
    input_tags: Dict[str, Any]

class ProcessStatusResponse(BaseModel):
    task_id: str
    status: str
    error: Optional[str] = None

class TagFilters(BaseModel):
    exact: Optional[Dict[str, Any]] = None
    range: Optional[Dict[str, List[Any]]] = None
    fuzzy: Optional[Dict[str, str]] = None

class SearchChunksRequest(BaseModel):
    query: str
    collection: Optional[str] = None
    top_k: int = 5
    fetch_k: int = 20
    dense_weight: float = 0.5
    sparse_weight: float = 0.5
    tag_filters: Optional[TagFilters] = None
    rerank: bool = False
    score_threshold: float = 0.0
    use_sparse: bool = False

class ChunkResult(BaseModel):
    id: str
    text: str
    score: float
    metadata: Dict[str, Any]

class SearchChunksResponse(BaseModel):
    results: List[ChunkResult]

class QueryRequest(BaseModel):
    query: str
    tag_filters: Optional[TagFilters] = None
    top_k: int = 5

class QueryResponse(BaseModel):
    answer: str
    sources: List[ChunkResult]

# --- Helper Functions ---
def extract_text_pdf(content: bytes) -> str:
    """Extract text from PDF using PyPDF2."""
    text = ""
    try:
        reader = PdfReader(io.BytesIO(content))
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    except Exception as e:
        print(f"Error extracting PDF text: {e}")
    return text

def chunk_text(text: str, chunk_size=500, overlap=50):
    """Split text into chunks of ~500 characters with overlap."""
    chunks = []
    if not text:
        return chunks
    step = chunk_size - overlap
    for i in range(0, len(text), step):
        chunk = text[i:i + chunk_size]
        if chunk:
            chunks.append(chunk)
    return chunks

async def get_embedding(text: str) -> List[float]:
    """Get embedding for a single text string."""
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {API_KEY}"
    }
    embeddings_url = f"{LLM_HOST_URL.rstrip('/')}/embeddings"
    async with httpx.AsyncClient(timeout=60.0) as client:
        payload = {
            "model": EMBEDDING_MODEL_NAME,
            "input": [text]
        }
        response = await client.post(embeddings_url, headers=headers, json=payload)
        if response.status_code != 200:
            raise Exception(f"Embedding API error: {response.status_code} - {response.text}")
        data = response.json()
        return data["data"][0]["embedding"]

async def background_process_file(file_id: int, task_id: str, input_tags: Dict[str, Any]):
    """Complete RAG workflow: extract, chunk, embed, store."""
    tasks_store[task_id]["status"] = "PROCESSING"
    try:
        file_data = uploaded_files_store.get(file_id)
        if not file_data:
            tasks_store[task_id]["status"] = "FAILED"
            tasks_store[task_id]["error"] = f"File ID {file_id} not found."
            return

        # Save to temp directory as requested
        file_path = TEMP_DIR / file_data["name"]
        with open(file_path, "wb") as f:
            f.write(file_data["content"])

        # Extract text
        ext = file_data["extension"]
        if ext == "pdf":
            text = extract_text_pdf(file_data["content"])
        elif ext in ["txt", "md"]:
            text = file_data["content"].decode("utf-8", errors="ignore")
        elif ext == "docx":
            from docx import Document
            doc = Document(io.BytesIO(file_data["content"]))
            text = "\n".join([para.text for para in doc.paragraphs])
        elif ext == "pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_data["content"]))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            text = "\n".join(text_runs)
        elif ext in ["xlsx", "xls"]:
            df = pd.read_excel(io.BytesIO(file_data["content"]))
            text = df.to_string()
        else:
            # Fallback for other files
            try:
                text = file_data["content"].decode("utf-8")
            except:
                text = file_data["content"].decode("latin-1", errors="ignore")

        if not text.strip():
            tasks_store[task_id]["status"] = "FAILED"
            tasks_store[task_id]["error"] = "No text extracted from file."
            return

        # Chunk text
        chunks = chunk_text(text)
        
        # Compute embeddings via OpenAI-compatible REST API
        embeddings = []
        batch_size = 16
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {API_KEY}"
        }
        # Construct embeddings endpoint URL
        embeddings_url = f"{LLM_HOST_URL.rstrip('/')}/embeddings"
        
        async with httpx.AsyncClient(timeout=60.0) as client:
            for i in range(0, len(chunks), batch_size):
                batch = chunks[i:i + batch_size]
                payload = {
                    "model": EMBEDDING_MODEL_NAME,
                    "input": batch
                }
                response = await client.post(embeddings_url, headers=headers, json=payload)
                if response.status_code != 200:
                    raise Exception(f"Embedding API error: {response.status_code} - {response.text}")
                
                data = response.json()
                # OpenAI format: response["data"] is a list of objects with "embedding" key
                embeddings.extend([item["embedding"] for item in data["data"]])
        
        # Prepare records for LanceDB
        raw_year = input_tags.get("academic_year", 0)
        try:
            if isinstance(raw_year, str):
                # Extract first 4 digits if string (e.g., "2023-2024")
                match = re.search(r"(\d{4})", raw_year)
                academic_year = int(match.group(1)) if match else 0
            else:
                academic_year = int(raw_year)
        except (ValueError, TypeError, AttributeError):
            academic_year = 0

        department = input_tags.get("department", "unknown")
        faculty = input_tags.get("faculty", "unknown")
        degree_level = input_tags.get("degree_level", "unknown")
        category = input_tags.get("category", "unknown")
        doc_title = file_data["name"]
        
        records = []
        for i, (chunk, vector) in enumerate(zip(chunks, embeddings)):
            records.append({
                "id": str(uuid.uuid4()),
                "text": chunk,
                "vector": vector,
                "academic_year": academic_year,
                "department": str(department),
                "faculty": str(faculty),
                "degree_level": str(degree_level),
                "category": str(category),
                "doc_title": str(doc_title),
                "chunk_index": i
            })
        
        # Store in LanceDB
        if TABLE_NAME not in db.list_tables().tables:
            db.create_table(TABLE_NAME, data=records)
        else:
            table = db.open_table(TABLE_NAME)
            table.add(records)
        
        tasks_store[task_id]["status"] = "COMPLETED"
    except Exception as e:
        tasks_store[task_id]["status"] = "FAILED"
        tasks_store[task_id]["error"] = str(e)

# --- Endpoints ---
@app.get("/health")
async def health_check():
    """Health check endpoint."""
    return {"status": "healthy"}

@app.post("/truncate")
async def truncate_database():
    """Drops the current table to truncate the database."""
    try:
        # Check if table exists. Note: list_tables() behavior can vary by version.
        tables = db.list_tables()
        if hasattr(tables, "tables"):
            table_list = tables.tables
        else:
            table_list = tables
            
        if TABLE_NAME in table_list:
            db.drop_table(TABLE_NAME)
            return {"status": "success", "message": f"Table {TABLE_NAME} truncated."}
        return {"status": "success", "message": f"Table {TABLE_NAME} did not exist."}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Truncation error: {str(e)}")

@app.post("/upload_file", response_model=FileMetadata)
async def upload_file(file: UploadFile = File(...)):
    """Accepts a file, returns file metadata."""
    global file_id_counter
    content = await file.read()
    file_id_counter += 1
    file_id = file_id_counter
    
    filename = file.filename
    extension = os.path.splitext(filename)[1].lower().lstrip('.')
    size = len(content)
    mime_type = file.content_type or "application/octet-stream"
    created_at = int(time.time())
    
    uploaded_files_store[file_id] = {
        "id": file_id,
        "name": filename,
        "content": content,
        "mime_type": mime_type,
        "extension": extension,
        "size": size,
        "created_at": created_at
    }
    
    return FileMetadata(
        id=file_id,
        name=filename,
        size=size,
        extension=extension,
        mime_type=mime_type,
        created_at=created_at
    )

@app.post("/process")
async def process_file(request: ProcessRequest, background_tasks: BackgroundTasks):
    """Starts async background RAG workflow."""
    if request.file_id not in uploaded_files_store:
        task_id = str(uuid.uuid4())
        tasks_store[task_id] = {"status": "FAILED", "error": f"File ID {request.file_id} not found."}
        return {"task_id": task_id}
    
    task_id = str(uuid.uuid4())
    tasks_store[task_id] = {"status": "PENDING", "error": None}
    
    background_tasks.add_task(background_process_file, request.file_id, task_id, request.input_tags)
    
    return {"task_id": task_id}

@app.get("/process/status/{task_id}", response_model=ProcessStatusResponse)
async def process_status(task_id: str):
    """Returns task status."""
    task = tasks_store.get(task_id)
    if not task:
        raise HTTPException(status_code=404, detail="Task ID not found")
    
    return ProcessStatusResponse(
        task_id=task_id,
        status=task["status"],
        error=task.get("error")
    )

@app.post("/search_chunks", response_model=SearchChunksResponse)
async def search_chunks(request: SearchChunksRequest):
    """Performs vector similarity search with metadata filtering."""
    if TABLE_NAME not in db.list_tables().tables:
        return SearchChunksResponse(results=[])
    
    table = db.open_table(TABLE_NAME)
    
    # Get query embedding
    try:
        query_vector = await get_embedding(request.query)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Embedding error: {str(e)}")
    
    # Translate tag_filters to LanceDB WHERE clause
    where_clauses = []
    if request.tag_filters:
        # Exact: any key-value pair
        if request.tag_filters.exact:
            for field, val in request.tag_filters.exact.items():
                if isinstance(val, str):
                    where_clauses.append(f"{field} = '{val}'")
                else:
                    where_clauses.append(f"{field} = {val}")
        
        # Range: any key with [min, max] list
        if request.tag_filters.range:
            for field, r in request.tag_filters.range.items():
                if isinstance(r, list) and len(r) == 2:
                    where_clauses.append(f"{field} BETWEEN {r[0]} AND {r[1]}")
        
        # Fuzzy: any key with substring
        if request.tag_filters.fuzzy:
            for field, val in request.tag_filters.fuzzy.items():
                where_clauses.append(f"{field} LIKE '%{val}%'")

    where_str = " AND ".join(where_clauses) if where_clauses else None
    
    # Perform vector search
    query_builder = table.search(query_vector)
    if where_str:
        query_builder = query_builder.where(where_str)
    
    candidates = query_builder.limit(request.top_k).to_list()
    
    results = []
    for cand in candidates:
        results.append(ChunkResult(
            id=cand["id"],
            text=cand["text"],
            score=float(cand.get("_distance", 0.0)),
            metadata={
                "academic_year": cand.get("academic_year"),
                "department": cand.get("department"),
                "faculty": cand.get("faculty"),
                "degree_level": cand.get("degree_level"),
                "category": cand.get("category"),
                "doc_title": cand.get("doc_title"),
                "chunk_index": cand.get("chunk_index")
            }
        ))
    
    return SearchChunksResponse(results=results)

@app.post("/query", response_model=QueryResponse)
async def query_rag(request: QueryRequest):
    """RAG endpoint: search then generate answer."""
    # 1. Search for relevant chunks
    search_req = SearchChunksRequest(
        query=request.query,
        tag_filters=request.tag_filters,
        top_k=request.top_k
    )
    search_resp = await search_chunks(search_req)
    
    if not search_resp.results:
        return QueryResponse(answer="I couldn't find any relevant information in the documents.", sources=[])
    
    # 2. Construct prompt
    context = "\n\n".join([f"Source {i+1}:\n{res.text}" for i, res in enumerate(search_resp.results)])
    system_prompt = "You are a helpful academic assistant. Answer the user's question based ONLY on the provided context. If the answer is not in the context, say you don't know."
    user_prompt = f"Context:\n{context}\n\nQuestion: {request.query}"
    
    # 3. Call LLM
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {API_KEY}"
    }
    chat_url = f"{LLM_HOST_URL.rstrip('/')}/chat/completions"
    
    async with httpx.AsyncClient(timeout=60.0) as client:
        payload = {
            "model": LLM_MODEL_NAME,
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            "temperature": 0.1
        }
        response = await client.post(chat_url, headers=headers, json=payload)
        if response.status_code != 200:
            raise HTTPException(status_code=500, detail=f"LLM API error: {response.status_code} - {response.text}")
        
        data = response.json()
        answer = data["choices"][0]["message"]["content"]
    
    return QueryResponse(answer=answer, sources=search_resp.results)

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host=RAG_HOST_PATH, port=RAG_PORT)
